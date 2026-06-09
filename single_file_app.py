from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import os
import uuid
import subprocess
import threading
import zipfile
import img2pdf
import pdfplumber
import pytesseract
import pikepdf
import html
from pdf2image import convert_from_path
from pypdf import PdfReader, PdfWriter
import pikepdf
import fitz
from docx import Document
import docx.shared
from PIL import Image as PILImage
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
import io
from deep_translator import GoogleTranslator
try:
    from weasyprint import HTML
except (ImportError, OSError):
    HTML = None

# ------------------------
# CONFIG
# ------------------------


app = Flask(__name__)
CORS(app)

import tempfile

UPLOAD_FOLDER = tempfile.gettempdir()
MAX_FILE_SIZE = 70 * 1024 * 1024
DELETE_DELAY = 60

from flask import send_from_directory

@app.route("/") # pyright: ignore[reportUndefinedVariable]
def home():
    return send_from_directory(".", "index.html")

@app.route("/ads.txt")
def ads_txt():
    return send_from_directory(".", "ads.txt")

@app.route("/robots.txt")
def robots_txt():
    return send_from_directory(".", "robots.txt")



os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# ------------------------
# UTILITIES
# ------------------------

def save_file(file):
    print(f"DEBUG: Saving file: {file.filename}")
    if file.content_length and file.content_length > MAX_FILE_SIZE:
        print("DEBUG: File exceeds size limit")
        raise Exception("File exceeds 70MB limit")

    filename = str(uuid.uuid4()) + "_" + file.filename
    path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(path)
    print(f"DEBUG: Saved to {path}")
    return path


def auto_delete(path):
    def delete():
        print(f"DEBUG: Deleting {path}")
        if os.path.exists(path):
            os.remove(path)
    threading.Timer(DELETE_DELAY, delete).start()

def clean_extracted_text(text):
    if not text:
         return ""
    replacements = {
        '\u2018': "'", '\u2019': "'",
        '\u201c': '"', '\u201d': '"',
        '\u2013': '-', '\u2014': '-',
        '\u2022': '-', '\u00B7': '-',
        'â€¢': '-', 'â€œ': '"', 'â€\x9d': '"', 'â€\x98': "'", 'â€\x99': "'", 'â€“': '-', 'â€”': '-',
        'œ': '"', '': '"', '˜': "'", '™': "'",
        'â€c': '"'
    }
    for k, v in replacements.items():
        text = text.replace(k, v)
    return text


# ------------------------
# MERGE
# ------------------------

@app.route("/merge", methods=["POST"])
def merge():
    try:
        files = request.files.getlist("files")
        writer = PdfWriter()
        paths = []

        for f in files:
            if not f.filename.lower().endswith('.pdf'):
                return jsonify({"error": f"File '{f.filename}' is not a PDF. Please select only PDF files for merging."}), 400
            path = save_file(f)
            paths.append(path)
            
            # Explicit Validation: Check for actual PDF byte signature
            with open(path, 'rb') as file_obj:
                header = file_obj.read(5)
            
            if header != b'%PDF-':
                converted = False
                try:
                    # Attempt text-to-PDF auto-conversion for fake/renamed PDFs
                    with open(path, 'r', encoding='utf-8', errors='ignore') as f_text:
                        content = f_text.read(500000) # Read up to 500KB of text
                    
                    if content.strip():
                        print(f"DEBUG: Auto-converting fake PDF '{f.filename}' to real PDF...")
                        new_path = path + "_converted.pdf"
                        
                        doc_pdf = SimpleDocTemplate(new_path, pagesize=letter)
                        styles = getSampleStyleSheet()
                        story = []
                        
                        for line in content.split('\n'):
                            if line.strip():
                                safe_line = html.escape(line.strip())
                                story.append(Paragraph(safe_line, styles["Normal"]))
                                
                        doc_pdf.build(story)
                        
                        auto_delete(path) # Cleanup fake PDF
                        path = new_path   # Swap in the new real PDF
                        paths[-1] = path  # Update paths list to delete the new file instead
                        converted = True
                except Exception as conv_e:
                    print(f"DEBUG: Auto-conversion failed: {conv_e}")
                
                if not converted:
                    err_msg = ""
                    if header.startswith(b'\xef\xbb\xbf') or header.startswith(b'<html') or header.startswith(b'<!DOC'):
                        err_msg = f"File '{f.filename}' is NOT a valid PDF document. It appears to be a plain text or HTML file saved with a .pdf extension, and auto-conversion failed."
                    else:
                        err_msg = f"File '{f.filename}' is not a valid PDF document (missing standard %PDF header)."
                    return jsonify({"error": err_msg}), 400

            try:
                # Attempt to read the PDF, even if it is slightly malformed (e.g. EOF missing)
                reader = PdfReader(path, strict=False)
                writer.append(reader)
            except Exception as first_error:
                try:
                    # PyPDF2 failed completely. Attempt a structural rebuild using pikepdf (QPDF engine)
                    print(f"DEBUG: PyPDF2 failed to read {f.filename}. Attempting pikepdf repair...")
                    repaired_path = path + "_repaired.pdf"
                    with pikepdf.open(path, allow_overwriting_input=False) as ppdf:
                        ppdf.save(repaired_path)
                    
                    # If repair successful, read the rebuilt file with PyPDF2
                    reader = PdfReader(repaired_path, strict=False)
                    writer.append(reader)
                    paths.append(repaired_path) # Mark the repaired temp file for cleanup
                except Exception as second_error:
                    # Provide a much cleaner and user-friendly error string
                    err_msg = str(first_error).replace("Stream has ended unexpectedly", "The file was extremely truncated or corrupted during download. Automatic repair failed.")
                    return jsonify({"error": f"Error merging '{f.filename}': {err_msg}"}), 400

        output = os.path.join(UPLOAD_FOLDER, str(uuid.uuid4()) + ".pdf")

        with open(output, "wb") as f:
            writer.write(f)

        for p in paths:
            auto_delete(p)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# SPLIT
# ------------------------

def parse_ranges(range_str, max_pages):
    pages = set()
    for part in range_str.split(','):
        part = part.strip()
        if not part: continue
        if '-' in part:
            try:
                start, end = part.split('-', 1)
                start = max(1, int(start.strip()))
                end = min(max_pages, int(end.strip()))
                if start <= end:
                    pages.update(range(start, end + 1))
            except ValueError:
                pass
        else:
            try:
                p = int(part)
                if 1 <= p <= max_pages:
                    pages.add(p)
            except ValueError:
                pass
    return sorted(list(pages))

@app.route("/split", methods=["POST"])
def split():
    try:
        import zipfile
        file_path = save_file(request.files["file"])
        reader = PdfReader(file_path)
        total_pages = len(reader.pages)
        
        split_mode = request.form.get('splitMode', 'extract')
        split_merge = request.form.get('splitMerge', 'false').lower() == 'true'
        
        output_files = []
        
        if split_mode == 'custom':
            range_str = request.form.get('splitRanges', '')
            pages_to_extract = parse_ranges(range_str, total_pages)
            if not pages_to_extract:
                return jsonify({"error": "No valid pages or ranges provided."}), 400
            
            if split_merge:
                writer = PdfWriter()
                for p in pages_to_extract:
                    writer.add_page(reader.pages[p - 1])
                out_path = os.path.join(UPLOAD_FOLDER, f"split_{uuid.uuid4()}.pdf")
                with open(out_path, "wb") as f_out:
                    writer.write(f_out)
                output_files.append(out_path)
            else:
                for p in pages_to_extract:
                    writer = PdfWriter()
                    writer.add_page(reader.pages[p - 1])
                    out_path = os.path.join(UPLOAD_FOLDER, f"page_{p}_{uuid.uuid4()}.pdf")
                    with open(out_path, "wb") as f_out:
                        writer.write(f_out)
                    output_files.append((out_path, f"page_{p}.pdf"))

        elif split_mode == 'fixed':
            try:
                chunk_size = int(request.form.get('splitFixedCount', 1))
            except ValueError:
                chunk_size = 1
            if chunk_size < 1: chunk_size = 1
            
            part = 1
            for i in range(0, total_pages, chunk_size):
                writer = PdfWriter()
                for j in range(i, min(i + chunk_size, total_pages)):
                    writer.add_page(reader.pages[j])
                out_path = os.path.join(UPLOAD_FOLDER, f"part_{part}_{uuid.uuid4()}.pdf")
                with open(out_path, "wb") as f_out:
                    writer.write(f_out)
                end_page = min(i + chunk_size, total_pages)
                output_files.append((out_path, f"pages_{i+1}_to_{end_page}.pdf"))
                part += 1
                
        elif split_mode == 'extract':
            for p in range(total_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[p])
                out_path = os.path.join(UPLOAD_FOLDER, f"page_{p+1}_{uuid.uuid4()}.pdf")
                with open(out_path, "wb") as f_out:
                    writer.write(f_out)
                output_files.append((out_path, f"page_{p+1}.pdf"))
                
        else:
            return jsonify({"error": "Unknown split mode."}), 400

        # Decide whether to return a single PDF or a ZIP
        if len(output_files) == 1 and isinstance(output_files[0], str):
            final_output = output_files[0]
            auto_delete(file_path)
            auto_delete(final_output)
            return send_file(final_output, as_attachment=True, download_name="split_result.pdf")
        else:
            zip_filename = f"split_files_{uuid.uuid4()}.zip"
            zip_filepath = os.path.join(UPLOAD_FOLDER, zip_filename)
            
            with zipfile.ZipFile(zip_filepath, 'w') as zipf:
                for item in output_files:
                    if isinstance(item, tuple):
                        path, name = item
                        zipf.write(path, name)
                        os.remove(path)
                    else:
                        zipf.write(item, os.path.basename(item))
                        os.remove(item)

            auto_delete(file_path)
            auto_delete(zip_filepath)
            return send_file(zip_filepath, as_attachment=True, download_name="split_files.zip")

    except Exception as e:
        print(f"Error in split: {e}")
        return jsonify({"error": str(e)}), 400


# ------------------------
# ORGANIZE
# ------------------------

@app.route("/organize", methods=["POST"])
def organize():
    try:
        file = save_file(request.files["file"])
        pages_str = request.form.get("pages", "").strip()
        output = file.replace(".pdf", "_organized.pdf")

        if not pages_str:
            return jsonify({"error": "No page order provided."}), 400

        # Parse requested pages (e.g. "4, 1, 2, 3" -> [4, 1, 2, 3])
        requested_pages = []
        for part in pages_str.split(','):
            part = part.strip()
            if not part: continue
            try:
                p = int(part)
                # Ensure it's not strictly 0 or negative
                if p > 0:
                    requested_pages.append(p)
            except ValueError:
                pass

        if not requested_pages:
             return jsonify({"error": "Invalid page numbers provided."}), 400

        reader = PdfReader(file)
        writer = PdfWriter()
        total_pages = len(reader.pages)

        # Build the new PDF in requested order
        # Also handles deletion implicitly (if a page is omitted, it's deleted)
        # And duplication (if a page is listed twice, it's duplicated)
        for p in requested_pages:
            # Check bounds
            if p <= total_pages:
                writer.add_page(reader.pages[p - 1])
            else:
                 return jsonify({"error": f"Requested page {p}, but document only has {total_pages} pages."}), 400

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        print(f"Error in organize: {e}")
        return jsonify({"error": str(e)}), 400


# ------------------------
# COMPRESS
# ------------------------

@app.route("/compress", methods=["POST"])
def compress():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_compressed.pdf")

        # Use pypdf for compression to avoid external dependencies like Ghostscript
        reader = PdfReader(file)
        writer = PdfWriter()
        
        # append and compress streams is much faster than compress_identical_objects for 10s timeouts
        writer.append(reader)
        for page in writer.pages:
            page.compress_content_streams()

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# WORD → PDF
# ------------------------

@app.route("/word-to-pdf", methods=["POST"])
@app.route("/word-to-pdf", methods=["POST"])
def word_to_pdf():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".docx", ".pdf").replace(".doc", ".pdf")

        # Setup ReportLab document
        doc_pdf = SimpleDocTemplate(output, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # 1. Extract Text from Paragraphs
        doc_word = Document(file)
        for para in doc_word.paragraphs:
            if para.text.strip():
                # Replace newlines to avoid reportlab errors
                clean_text = para.text.replace("\n", "<br/>")
                story.append(Paragraph(clean_text, styles["Normal"]))
                story.append(Spacer(1, 12))

        # 2. Extract Text from Tables (Basic)
        if doc_word.tables:
            story.append(Paragraph("<b>Tables Content:</b>", styles["Heading2"]))
            for table in doc_word.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    # reportlab Table
                    t = Table(table_data)
                    t.setStyle(TableStyle([
                        ('GRID', (0,0), (-1,-1), 1, colors.black),
                        ('FONTNAME', (0,0), (-1,-1), 'Helvetica'),
                        ('FONTSIZE', (0,0), (-1,-1), 10),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 12))

        # 3. Extract Images from the DOCX (it is a ZIP file)
        # We can't easily place them strictly in-line with text without complex XML parsing,
        # so we append them after the text.
        with zipfile.ZipFile(file, 'r') as z:
            # Find all files in word/media/
            media_files = [f for f in z.namelist() if f.startswith("word/media/")]
            if media_files:
                story.append(PageBreak())
                story.append(Paragraph("<b>Images:</b>", styles["Heading2"]))
                
                for media_path in sorted(media_files):
                    # Extract image to temp file
                    img_data = z.read(media_path)
                    img_name = os.path.basename(media_path)
                    temp_img_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4()}_{img_name}")
                    
                    with open(temp_img_path, "wb") as f_img:
                        f_img.write(img_data)
                    
                    try:
                        # Add to PDF
                        # Resize if too big? SimpleDocTemplate handles page fit mostly, but let's be safe
                        img = Image(temp_img_path)
                        
                        # Basic scaling logic to fit page width (approx 6 inches = 432 pts)
                        max_width = 450
                        if img.drawWidth > max_width:
                            ratio = max_width / img.drawWidth
                            img.drawWidth = max_width
                            img.drawHeight = img.drawHeight * ratio
                            
                        story.append(img)
                        story.append(Spacer(1, 12))
                        
                        # Clean up temp image immediately found? No, reportlab needs it until build.
                        # We'll rely on auto_delete or deleting after build.
                        # Actually reportlab needs the file to exist during .build()
                        # We will delete them after the build.
                    except Exception as img_err:
                        print(f"DEBUG: Could not add image {media_path}: {img_err}")

        # Build PDF
        doc_pdf.build(story)

        # Cleanup temp images - we need to iterate directory or just rely on OS/restart
        # For this logic, let's just leave them for the generic cleanup or direct remove if we tracked them.
        # Tracked deletion:
        # (Optimally we'd list them, but uuid helps collision. We'll leave them for the 1hr cleanup func)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        print(f"DEBUG: Word to PDF error: {e}")
        return jsonify({"error": str(e)}), 400


# ------------------------
# PDF → WORD
# ------------------------

@app.route("/pdf-to-word", methods=["POST"])
def pdf_to_word():
    try:
        file = save_file(request.files["file"])
        doc = Document()
        text_found = False

        # 1. Try to extract text
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    doc.add_paragraph(text)
                    text_found = True

        # 2. If no text found, assume it's a scanned PDF and extract images
        if not text_found:
            print("DEBUG: No text found in PDF, attempting image extraction...")
            reader = PdfReader(file)
            for page in reader.pages:
                for image in page.images:
                    # image.data contains the bytes
                    image_stream = io.BytesIO(image.data)
                    try:
                        doc.add_picture(image_stream, width=docx.shared.Inches(6))
                    except Exception as img_err:
                         print(f"DEBUG: Failed to add image: {img_err}")
                    doc.add_paragraph("") # Add spacing

        output = file.replace(".pdf", ".docx")
        doc.save(output)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# JPG → PDF
# ------------------------

@app.route("/jpg-to-pdf", methods=["POST"])
def jpg_to_pdf():
    try:
        files = request.files.getlist("files")
        paths = [save_file(f) for f in files]
        
        pagesize = request.form.get('pagesize', 'fit')

        output = os.path.join(UPLOAD_FOLDER, str(uuid.uuid4()) + ".pdf")

        layout_fun = None
        if pagesize.lower() == 'a4':
            a4inpt = (img2pdf.in_to_pt(8.27), img2pdf.in_to_pt(11.69))
            layout_fun = img2pdf.get_layout_fun(pagesize=a4inpt, fit=img2pdf.FitMode.into)
        elif pagesize.lower() == 'a3':
            a3inpt = (img2pdf.in_to_pt(11.69), img2pdf.in_to_pt(16.54))
            layout_fun = img2pdf.get_layout_fun(pagesize=a3inpt, fit=img2pdf.FitMode.into)

        with open(output, "wb") as f:
            if layout_fun:
                f.write(img2pdf.convert(paths, rotation=img2pdf.Rotation.ifvalid, layout_fun=layout_fun))
            else:
                f.write(img2pdf.convert(paths, rotation=img2pdf.Rotation.ifvalid))

        for p in paths:
            auto_delete(p)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# PDF → JPG
# ------------------------

@app.route("/pdf-to-jpg", methods=["POST"])
def pdf_to_jpg():
    try:
        files = request.files.getlist("files")
        if not files and "file" in request.files:
            files = [request.files["file"]]

        if not files:
             return jsonify({"error": "No file uploaded"}), 400
             
        all_images = []
        temp_files = []
        errors = []

        print(f"DEBUG: Processing {len(files)} files for PDF-to-JPG")

        for file_obj in files:
            file_path = save_file(file_obj)
            temp_files.append(file_path)
            
            # Create a safe base name
            import re
            safe_name = re.sub(r'[^a-zA-Z0-9_\-]', '_', file_obj.filename)
            base_name = os.path.splitext(safe_name)[0] or "image"
            
            # Try PyMuPDF
            success_fitz = False
            try:
                doc = fitz.open(file_path)
                if len(doc) == 0:
                    errors.append(f"PyMuPDF: Document has 0 pages ({file_obj.filename})")
                else:
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=200) # Slightly lower DPI to save memory on Vercel
                        img_filename = f"{base_name}_page_{i+1}.jpg"
                        img_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_{img_filename}")
                        pix.save(img_path)
                        if os.path.exists(img_path) and os.path.getsize(img_path) > 0:
                            all_images.append((img_path, img_filename))
                    success_fitz = True
                doc.close()
            except Exception as e:
                errors.append(f"PyMuPDF failed ({type(e).__name__}): {str(e)}")
            
            # Fallback to pypdfium2
            if not success_fitz:
                try:
                    import pypdfium2 as pdfium
                    pdf = pdfium.PdfDocument(file_path)
                    if len(pdf) == 0:
                        errors.append(f"pypdfium2: Document has 0 pages ({file_obj.filename})")
                    else:
                        for i in range(len(pdf)):
                            page = pdf[i]
                            # scale=2 is ~144 dpi, scale=3 is ~216 dpi
                            bitmap = page.render(scale=2.5)
                            pil_image = bitmap.to_pil()
                            img_filename = f"{base_name}_page_{i+1}.jpg"
                            img_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_{img_filename}")
                            pil_image.save(img_path, "JPEG", quality=90)
                            if os.path.exists(img_path):
                                all_images.append((img_path, img_filename))
                        success_pypdfium2 = True
                    pdf.close()
                except Exception as e2:
                    errors.append(f"pypdfium2 failed ({type(e2).__name__}): {str(e2)}")

        # Check if we got ANY images
        if not all_images:
            error_message = "No images could be extracted from the provided PDF(s)."
            if errors:
                error_message += " Details: " + " | ".join(errors)
            return jsonify({"error": error_message}), 400

        # Create output
        if len(all_images) == 1:
            output = all_images[0][0]
        else:
            output = os.path.join(UPLOAD_FOLDER, f"converted_images_{uuid.uuid4()}.zip")
            with zipfile.ZipFile(output, 'w') as zipf:
                for img_path, img_name in all_images:
                    zipf.write(img_path, img_name)
                    os.remove(img_path)

        for p in temp_files:
            auto_delete(p)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        print(f"Error in pdf_to_jpg: {e}")
        return jsonify({"error": str(e)}), 400


# ------------------------
# PROTECT
# ------------------------

@app.route("/protect", methods=["POST"])
def protect():
    try:
        file = save_file(request.files["file"])
        password = request.form.get("password")

        reader = PdfReader(file)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)

        output = file.replace(".pdf", "_protected.pdf")

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# UNLOCK
# ------------------------

@app.route("/unlock", methods=["POST"])
def unlock():
    try:
        file = save_file(request.files["file"])
        password = request.form.get("password")
        output = file.replace(".pdf", "_unlocked.pdf")

        # Use pikepdf for reliable decryption
        try:
            with pikepdf.open(file, password=password) as pdf:
                pdf.save(output)
        except pikepdf.PasswordError:
             return jsonify({"error": "Incorrect password"}), 400

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# WATERMARK
# ------------------------

@app.route("/watermark", methods=["POST"])
def watermark():
    try:
        file = save_file(request.files["file"])
        watermark_text = request.form.get("watermarkText", "Watermark")
        password = request.form.get("password")
        output = file.replace(".pdf", "_watermarked.pdf")

        if password:
            # Use pikepdf for reliable decryption first
            try:
                unlocked_file = file.replace(".pdf", "_temp_unlocked.pdf")
                with pikepdf.open(file, password=password) as pdf:
                    pdf.save(unlocked_file)
                # Switch to using the unlocked file for processing
                auto_delete(file) # Delete original locked file from temp
                file = unlocked_file # Point 'file' to the new unlocked version
            except pikepdf.PasswordError:
                return jsonify({"error": "Incorrect password"}), 400

        reader = PdfReader(file)
        # No need to decrypt with pypdf anymore as we trust pikepdf did it
        
        writer = PdfWriter()

        # Create a temporary watermark for each page logic (optimized by reuse if sizes match, but simple loop here)
        # Actually, if pages differ in size, we need dynamic watermarks.
        # But for simplicity, let's create one based on the first page, or generate per page?
        # Generating per page is safer but slower. Let's do a smart approach:
        # Calculate center based on page.mediabox

        for page in reader.pages:
            # Get page dimensions
            # mediabox is [x, y, width, height] usually
            page_width = float(page.mediabox.width)
            page_height = float(page.mediabox.height)

            # Create a localized watermark PDF in memory
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=(page_width, page_height))
            
            # Setup Grid/Font
            c.setFont("Helvetica", 50)
            c.setFillColorRGB(0.5, 0.5, 0.5, 0.5)
            c.saveState()
            
            # Move to center
            c.translate(page_width / 2, page_height / 2)
            c.rotate(45)
            c.drawCentredString(0, 0, watermark_text)
            
            c.restoreState()
            c.save()

            packet.seek(0)
            watermark_reader = PdfReader(packet)
            watermark_page = watermark_reader.pages[0]

            # Merge
            page.merge_page(watermark_page)
            writer.add_page(page)

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# OCR
# ------------------------

@app.route("/ocr", methods=["POST"])
def ocr():
    print("OCR Request Received")
    try:
        # Check for Tesseract in common paths
        tesseract_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            os.path.expanduser(r"~\AppData\Local\Tesseract-OCR\tesseract.exe")
        ]
        
        found_tesseract = False
        for path in tesseract_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                print(f"Found Tesseract at: {path}")
                found_tesseract = True
                break
        
        if not found_tesseract:
            print("Tesseract not found in common paths, checking PATH...")
            # If not found, trust PATH, but might fail later.

        file = save_file(request.files["file"])
        print(f"File saved: {file}")
        
        lang = request.form.get("lang", "eng")
        print(f"Language selected: {lang}")

        # Use PyMuPDF instead of pdf2image to avoid Poppler dependency
        doc = fitz.open(file)
        
        output = file.rsplit(".", 1)[0] + ".txt"
        print(f"Output file will be: {output}")

        # Determine tessdata path
        # Check if local tessdata exists
        local_tessdata = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tessdata")
        if os.path.exists(local_tessdata):
             tessdata_path = local_tessdata
             os.environ["TESSDATA_PREFIX"] = tessdata_path
             print(f"Setting TESSDATA_PREFIX to local: {tessdata_path}")
             
             # Check if language file exists
             lang_file = os.path.join(tessdata_path, f"{lang}.traineddata")
             if not os.path.exists(lang_file):
                print(f"Missing language file: {lang_file}")
                return jsonify({
                    "error": f"The '{lang}' language data is missing in local folder using path: {lang_file}. Please download '{lang}.traineddata'."
                }), 400

        else:
            # Fallback to system tessdata if tesseract_cmd is absolute
            cmd = pytesseract.pytesseract.tesseract_cmd
            if os.path.isabs(cmd):
                tesseract_folder = os.path.dirname(cmd)
                tessdata_path = os.path.join(tesseract_folder, "tessdata")
                os.environ["TESSDATA_PREFIX"] = tessdata_path
                print(f"Setting TESSDATA_PREFIX to system: {tessdata_path}")
                
                # Check if language file exists
                lang_file = os.path.join(tessdata_path, f"{lang}.traineddata")
                if not os.path.exists(lang_file):
                    print(f"Missing language file: {lang_file}")
                    return jsonify({
                        "error": f"The '{lang}' language data is missing. Please download '{lang}.traineddata' from https://github.com/tesseract-ocr/tessdata and place it in '{tessdata_path}'."
                    }), 400
            else:
                 print("Tesseract command is not an absolute path, skipping TESSDATA_PREFIX setup.")

        with open(output, "w", encoding="utf-8-sig") as f:
            for page_num in range(len(doc)):
                print(f"Processing page {page_num + 1}/{len(doc)}")
                try:
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(dpi=300) # High DPI for OCR
                    
                    # Convert to PIL Image
                    mode = "RGBA" if pix.alpha else "RGB"
                    img = PILImage.frombytes(mode, [pix.width, pix.height], pix.samples)
                    
                    # Run OCR
                    text = pytesseract.image_to_string(img, lang=lang)
                    f.write(text + "\n\n")
                    print(f"Page {page_num + 1} processed.")
                except Exception as e:
                    print(f"Error processing page {page_num + 1}: {str(e)}")
                    if "tesseract is not installed" in str(e):
                        raise Exception("Tesseract OCR is not installed. Please download and install it from: https://github.com/UB-Mannheim/tesseract/wiki")
                    raise e

        # Check if file was written
        if not os.path.exists(output):
             print("Output file was not created!")
             raise Exception("Failed to create output file.")

        print("OCR Completed. cleaning up input file.")
        auto_delete(file)
        # auto_delete(output) # Don't delete output yet, let send_file handle it or rely on auto_delete garbage collection if implemented differently. 
        # Actually current implementation of auto_delete deletes it after a delay.
        auto_delete(output)

        print("Sending file...")
        return send_file(output, as_attachment=True)
    except Exception as e:
        print(f"OCR FAILED: {str(e)}")
        # traceback.print_exc()
        return jsonify({"error": str(e)}), 400

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# PPT TO PDF
# ------------------------

@app.route("/ppt-to-pdf", methods=["POST"])
def ppt_to_pdf():
    try:
        file = save_file(request.files["file"])
        output = file.rsplit(".", 1)[0] + ".pdf"
        
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader

        prs = Presentation(file)
        # Create PDF with same size as first slide (assuming uniform)
        # pptx sizes are in EMU (914400 per inch). Reportlab uses points (72 per inch).
        # Factor: 914400 / 72 = 12700
        
        width = prs.slide_width / 12700
        height = prs.slide_height / 12700
        
        c = canvas.Canvas(output, pagesize=(width, height))
        
        for slide in prs.slides:
            # We can't easily reproduce the exact look (backgrounds, complex shapes)
            # But we can try to extract text and images and place them roughly correctly.
            
            for shape in slide.shapes:
                # Convert EMU to points
                x = shape.left / 12700
                # reportlab (0,0) is bottom-left, pptx (0,0) is top-left
                y = height - (shape.top / 12700) - (shape.height / 12700)
                w = shape.width / 12700
                h = shape.height / 12700
                
                if shape.shape_type == 13: # PICTURE
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        # Draw image
                        img_reader = ImageReader(io.BytesIO(image_bytes))
                        c.drawImage(img_reader, x, y, width=w, height=h)
                    except Exception:
                        pass
                
                if shape.has_text_frame:
                    # Draw text
                    # Simplified: just draw all text in the box
                    text = shape.text
                    if text.strip():
                        # styling is hard, let's just draw standard font
                        text_obj = c.beginText(x, y + h - 12) # Start from top of box
                        text_obj.setFont("Helvetica", 10)
                        # Split by lines
                        for line in text.split('\n'):
                            text_obj.textLine(line)
                        c.drawText(text_obj)
            
            c.showPage()
            
        c.save()

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# EXCEL TO PDF
# ------------------------

@app.route("/excel-to-pdf", methods=["POST"])
def excel_to_pdf():
    try:
        file = save_file(request.files["file"])
        output = file.rsplit(".", 1)[0] + ".pdf"
        
        from openpyxl import load_workbook
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle

        wb = load_workbook(file, data_only=True)
        ws = wb.active
        
        doc = SimpleDocTemplate(output, pagesize=landscape(letter))
        elements = []
        
        data = []
        for row in ws.iter_rows(values_only=True):
            # Replace None with empty string to avoid reportlab errors
            row_data = [str(cell) if cell is not None else "" for cell in row]
            data.append(row_data)
            
        if data:
            # Create Table
            t = Table(data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            elements.append(t)
            
        doc.build(elements)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# GROUP 1: PAGE MANIPULATION
# ------------------------

# ------------------------
# ROTATE PDF
# ------------------------
@app.route("/rotate", methods=["POST"])
def rotate():
    try:
        file = save_file(request.files["file"])
        angle = int(request.form.get("angle", 90))
        output = file.replace(".pdf", f"_rotated_{angle}.pdf")

        reader = PdfReader(file)
        writer = PdfWriter()

        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# REMOVE PAGES
# ------------------------
@app.route("/remove-pages", methods=["POST"])
def remove_pages():
    try:
        file = save_file(request.files["file"])
        pages_to_remove = request.form.get("pages") # "1,3,5" or "1-3"
        output = file.replace(".pdf", "_removed.pdf")

        reader = PdfReader(file)
        writer = PdfWriter()
        
        # Parse pages
        total_pages = len(reader.pages)
        to_remove = set()
        
        for part in pages_to_remove.split(','):
            if '-' in part:
                start, end = map(int, part.split('-'))
                to_remove.update(range(start, end + 1))
            else:
                to_remove.add(int(part))

        for i in range(total_pages):
            if (i + 1) not in to_remove:
                writer.add_page(reader.pages[i])

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# ADD PAGE NUMBERS
# ------------------------
@app.route("/add-page-numbers", methods=["POST"])
def add_page_numbers():
    try:
        file = save_file(request.files["file"])
        position = request.form.get("position", "bottom-center")
        output = file.replace(".pdf", "_numbered.pdf")

        reader = PdfReader(file)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            # Create a watermark PDF with the page number
            packet = io.BytesIO()
            # Use page size from the page itself
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            
            c = canvas.Canvas(packet, pagesize=(width, height))
            c.setFont("Helvetica", 12)
            
            text = f"Page {i + 1}"
            text_width = c.stringWidth(text, "Helvetica", 12)
            
            x, y = 0, 0
            margin = 20
            
            if "bottom" in position:
                y = margin
            elif "top" in position:
                y = height - margin
                
            if "center" in position:
                x = (width - text_width) / 2
            elif "left" in position:
                x = margin
            elif "right" in position:
                x = width - margin - text_width
                
            c.drawString(x, y, text)
            c.save()
            
            packet.seek(0)
            number_pdf = PdfReader(packet)
            page.merge_page(number_pdf.pages[0])
            writer.add_page(page)

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# CROP PDF
# ------------------------
@app.route("/crop", methods=["POST"])
def crop():
    try:
        # Expected inputs: margin_left, margin_right, margin_top, margin_bottom (in points or inches? assume points for now or percentage?)
        # Let's simple crop: "Remove whitespace" or "Crop specific margins"
        # For simplicity in this first pass, let's just do a symmetric crop or simple coordinate crop.
        # Better yet: "Crop 1 inch from all sides" ?
        # Let's accept: left, right, top, bottom offsets in points.
        
        file = save_file(request.files["file"])
        # Defaults to 0
        left = float(request.form.get("left", 0))
        right = float(request.form.get("right", 0))
        top = float(request.form.get("top", 0))
        bottom = float(request.form.get("bottom", 0))
        
        output = file.replace(".pdf", "_cropped.pdf")

        reader = PdfReader(file)
        writer = PdfWriter()
        
        for page in reader.pages:
            # page.mediabox is [llx, lly, urx, ury]
            # ll = lower-left, ur = upper-right
            
            current_llx = float(page.mediabox.left)
            current_lly = float(page.mediabox.bottom)
            current_urx = float(page.mediabox.right)
            current_ury = float(page.mediabox.top)
            
            new_llx = current_llx + left
            new_lly = current_lly + bottom
            new_urx = current_urx - right
            new_ury = current_ury - top
            
            # Validation
            if new_llx >= new_urx or new_lly >= new_ury:
                # Invalid crop, skip or fail? 
                # Let's just keep original if crop invalid
                writer.add_page(page)
                continue
                
            page.mediabox.lower_left = (new_llx, new_lly)
            page.mediabox.upper_right = (new_urx, new_ury)
            writer.add_page(page)

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# GROUP 2: CONVERSION & OCR
# ------------------------

# ------------------------
# WEB TO PDF
# ------------------------
@app.route("/web-to-pdf", methods=["POST"])
def web_to_pdf():
    try:
        url = request.form.get("url")
        if not url:
            return jsonify({"error": "No URL provided"}), 400
            
        output = os.path.join(UPLOAD_FOLDER, str(uuid.uuid4()) + ".pdf")
        
        success = False
        error_msg = ""
        
        # 1. Try Headless Browser (Chrome/Edge) for best fidelity
        try:
            browser_paths = [
                r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
                r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
                r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
            ]
            
            browser_exe = None
            for p in browser_paths:
                if os.path.exists(p):
                    browser_exe = p
                    break
            
            if browser_exe:
                print(f"Converting URL with Browser: {browser_exe}")
                # Use --no-pdf-header-footer to avoid url/date in margins if possible (Chrome supports it, Edge likely too)
                cmd = [
                    browser_exe,
                    "--headless",
                    "--disable-gpu",
                    "--no-pdf-header-footer",
                    f"--print-to-pdf={os.path.abspath(output)}",
                    url
                ]
                # Timeout to prevent hanging
                subprocess.run(cmd, check=True, timeout=30)
                
                if os.path.exists(output) and os.path.getsize(output) > 0:
                    success = True
            else:
                error_msg += "No supported browser found for high-fidelity conversion. "

        except Exception as br_err:
             print(f"Browser conversion failed: {br_err}")
             error_msg += f"Browser: {str(br_err)} | "

        # 2. Use WeasyPrint (Alternative)
        if not success:
            try:
                if HTML is None:
                    raise ImportError("WeasyPrint not installed")
                print(f"Converting URL with WeasyPrint: {url}")
                HTML(url).write_pdf(output)
                success = True
            except Exception as wp_err:
                print(f"WeasyPrint failed: {wp_err}")
                error_msg += f"WeasyPrint: {str(wp_err)} | "

        # 3. Fallback to xhtml2pdf
        if not success:
            try:
                print(f"Converting URL with xhtml2pdf fallback: {url}")
                import requests
                from xhtml2pdf import pisa
                
                # Fetch content
                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
                response = requests.get(url, headers=headers, timeout=15)
                if response.status_code != 200:
                    raise Exception(f"Failed to fetch URL: {response.status_code}")
                
                with open(output, "wb") as f:
                    pisa_status = pisa.CreatePDF(response.text, dest=f)
                
                if pisa_status.err:
                    raise Exception("xhtml2pdf conversion error")
                
                success = True
            except Exception as xp_err:
                print(f"xhtml2pdf failed: {xp_err}")
                error_msg += f"xhtml2pdf: {str(xp_err)}"
        
        if not success:
             return jsonify({"error": f"Conversion failed: {error_msg}"}), 500

        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# PDF TO PDF/A
# ------------------------
@app.route("/pdf-to-pdfa", methods=["POST"])
def pdf_to_pdfa():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_pdfa.pdf")

        # Try Ghostscript
        gs_cmd = None
        # Check specific paths or PATH
        possible_cmds = ["gswin64c", "gs", "ghostscript"]
        for cmd in possible_cmds:
            from shutil import which
            if which(cmd):
                gs_cmd = cmd
                break
        
        if gs_cmd:
            # GS command for PDF/A-2b
            # gs -dPDFA -dBATCH -dNOPAUSE -sProcessColorModel=DeviceCMYK -sDEVICE=pdfwrite -sPDFACompatibilityPolicy=1 -sOutputFile=output.pdf input.pdf
            cmd = [
                gs_cmd,
                "-dPDFA",
                "-dBATCH",
                "-dNOPAUSE",
                "-sProcessColorModel=DeviceRGB", # Use RGB for screen
                "-sDEVICE=pdfwrite",
                "-sPDFACompatibilityPolicy=1",
                f"-sOutputFile={output}",
                file
            ]
            subprocess.run(cmd, check=True)
            
        else:
            # Fallback: Just linearized (Fast Web View) using pikepdf
            try:
                import pikepdf
                with pikepdf.open(file) as pdf:
                    pdf.save(output, linearize=True)
                print("Ghostscript not found, falling back to Linearized PDF (Fast Web View)")
            except ImportError:
                 return jsonify({"error": "Ghostscript not found and pikepdf missing."}), 500

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# PDF TO EXCEL (OCR)
# ------------------------
@app.route("/pdf-to-excel", methods=["POST"])
def pdf_to_excel():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", ".xlsx")
        
        # Strategy:
        # 1. Try pdfplumber for tables
        # 2. If no tables/text, use OCR (tesseract) -> then put text in cells
        
        import pdfplumber
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Data"
        
        row_idx = 1
        
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                # Extract Tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            # Clean None
                            clean_row = [cell if cell else "" for cell in row]
                            ws.append(clean_row)
                            row_idx += 1
                        ws.append([]) # Empty row between tables
                        row_idx += 1
                else:
                    # Fallback to text extraction if no tables
                    text = page.extract_text()
                    if text:
                        for line in text.split('\n'):
                            ws.append([line])
                            row_idx += 1
                    else:
                        # Scanned? OCR time
                        # We recycle the OCR route logic or just call pytesseract
                        # For now, let's just use images extraction
                        print("No text found, attempting OCR for Excel...")
                        # This is expensive, might timeout.
                        # Simple implementation: Image -> String -> Cell
                        pass

        wb.save(output)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400

@app.route("/pdf-to-ppt", methods=["POST"])
def pdf_to_ppt():
    try:
        file = save_file(request.files["file"])
        output = file.rsplit(".", 1)[0] + ".pptx"
        
        from pptx import Presentation
        from pptx.util import Inches
        import fitz # PyMuPDF

        # Convert PDF to images using PyMuPDF (no Poppler needed)
        doc = fitz.open(file)
        
        prs = Presentation()
        # Use a blank layout
        blank_slide_layout = prs.slide_layouts[6]
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=150) # Good quality for PPT
            
            # Save temp image
            image_path = f"{file}_page_{page_num}.png"
            pix.save(image_path)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide, fitting the height
            # python-pptx default slide height is 7.5 inches = 6858000 EMU
            # We can use the image aspect ratio to center it.
            
            pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0), height=prs.slide_height)
            
            # Center horizontally if needed
            if pic.width < prs.slide_width:
                pic.left = int((prs.slide_width - pic.width) / 2)
            
            auto_delete(image_path)

        prs.save(output)

        auto_delete(file)
        auto_delete(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# GROUP 3: CONTENT EDITING & SECURITY
# ------------------------

# ------------------------
# REPAIR PDF
# ------------------------
@app.route("/repair", methods=["POST"])
def repair():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_repaired.pdf")
     
        # Use pikepdf to repair (it automatically fixes many issues on open)
        import pikepdf
        try:
            with pikepdf.open(file) as pdf:
                pdf.save(output)
        except Exception as e:
             # Fallback to qpdf via subprocess if installed? 
             # Or try pypdf logic
             print(f"Pikepdf repair failed: {e}, trying pypdf reconstruction...")
             reader = PdfReader(file)
             writer = PdfWriter()
             for page in reader.pages:
                 writer.add_page(page)
             with open(output, "wb") as f:
                 writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# SIGN PDF
# ------------------------
@app.route("/sign", methods=["POST"])
def sign():
    try:
        file = save_file(request.files["file"])
        signature_file = save_file(request.files["signature"])
        output = file.replace(".pdf", "_signed.pdf")
        
        # Coords (default to bottom right)
        x = float(request.form.get("x", 400))
        y = float(request.form.get("y", 50))
        scale = float(request.form.get("scale", 0.3)) # Scale factor for image
        page_num = int(request.form.get("page", 1)) - 1
        
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        
        reader = PdfReader(file)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            if i == page_num:
                # Create signature overlay
                packet = io.BytesIO()
                # Get page size
                # Handle varying page sizes
                try:
                    w = float(page.mediabox.width)
                    h = float(page.mediabox.height)
                except:
                    w, h = 595, 842 # Default A4
                    
                c = canvas.Canvas(packet, pagesize=(w, h))
                
                # Draw image
                img = ImageReader(signature_file)
                iw, ih = img.getSize()
                # Default width 150pts
                target_w = 150 * scale
                target_h = (ih / iw) * target_w
                
                c.drawImage(img, x, y, width=target_w, height=target_h, mask='auto')
                c.save()
                
                packet.seek(0)
                overlay = PdfReader(packet)
                page.merge_page(overlay.pages[0])
            
            writer.add_page(page)

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(signature_file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# REDACT PDF
# ------------------------
@app.route("/redact", methods=["POST"])
def redact():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_redacted.pdf")
        
        # Input: text to redact
        redact_text = request.form.get("text")
        
        doc = fitz.open(file)
        
        # Search and redact
        for page in doc:
            if redact_text:
                areas = page.search_for(redact_text)
                for area in areas:
                    page.add_redact_annot(area, fill=(0, 0, 0)) # Black box
                page.apply_redactions()
        
        doc.save(output)
        doc.close()

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# PDF FILLER (FORM)
# ------------------------
@app.route("/fill-form", methods=["POST"])
def fill_form():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_filled.pdf")
        
        # User provides JSON mapping
        import json
        fields_str = request.form.get("fields", "").strip()
        if not fields_str:
             fields_str = "{}"
        try:
             fields_data = json.loads(fields_str)
        except Exception as e:
             return jsonify({"error": "Invalid fields JSON format. Please ensure it is valid JSON like {\"Name\": \"John\"}"}), 400
        
        reader = PdfReader(file)
        writer = PdfWriter()
        
        writer.append(reader)
        
        if reader.is_encrypted:
             writer.encrypt(reader.password if reader.password else "")

        # Try updating fields
        # Note: This only works for AcroForms, not XFA
        writer.update_page_form_field_values(writer.pages[0], fields_data)
        
        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# EDIT PDF (Overlay Text)
# ------------------------
@app.route("/edit-pdf", methods=["POST"])
def edit_pdf():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_edited.pdf")
        
        # Text, X, Y, Page
        text = request.form.get("text", "")
        # Default coords
        x = float(request.form.get("x", 100))
        y = float(request.form.get("y", 100))
        page_num = int(request.form.get("page", 1)) - 1
        font_size = int(request.form.get("fontSize", 12))
        
        reader = PdfReader(file)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            if i == page_num:
                packet = io.BytesIO()
                try:
                    w = float(page.mediabox.width)
                    h = float(page.mediabox.height)
                except:
                     w, h = 595, 842
                     
                c = canvas.Canvas(packet, pagesize=(w, h))
                c.setFont("Helvetica", font_size)
                c.drawString(x, y, text)
                c.save()
                
                packet.seek(0)
                overlay = PdfReader(packet)
                page.merge_page(overlay.pages[0])
            
            writer.add_page(page)

        with open(output, "wb") as f:
            writer.write(f)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# ALIASES / IMPROVEMENTS
# ------------------------
@app.route("/scan-to-pdf", methods=["POST"])
def scan_to_pdf():
    # Alias for JPG to PDF (frontend should send files)
    # We call the existing function logic?
    # Flask routes wrap functions. We can just call logic if extracted, 
    # but here functions are views.
    # We can redirect or just copy logic. 
    # Easier to just call the function if it doesn't rely on request.method checks that fail.
    # jpg_to_pdf uses request.files.getlist("files")
    return jpg_to_pdf()


# ------------------------
# GROUP 4: ADVANCED & AI
# ------------------------

# ------------------------
# COMPARE PDF
# ------------------------
@app.route("/compare", methods=["POST"])
def compare():
    try:
        # Expect two files
        files = request.files.getlist("files")
        if len(files) < 2:
            return jsonify({"error": "Please upload two PDF files to compare."}), 400
            
        file1 = save_file(files[0])
        file2 = save_file(files[1])
        output = os.path.join(UPLOAD_FOLDER, f"comparison_{uuid.uuid4()}.html")
        
        import difflib
        
        # Extract text
        def get_text(path):
            doc = fitz.open(path)
            text = ""
            for page in doc:
                text += clean_extracted_text(page.get_text())
            return text
            
        text1 = get_text(file1).splitlines()
        text2 = get_text(file2).splitlines()
        
        d = difflib.HtmlDiff()
        html = d.make_file(text1, text2, fromdesc=files[0].filename, todesc=files[1].filename)
        
        with open(output, "w", encoding="utf-8-sig") as f:
            f.write(html)
            
        auto_delete(file1)
        auto_delete(file2)
        auto_delete(output)
        
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# TRANSLATE PDF
# ------------------------
@app.route("/translate", methods=["POST"])
def translate_pdf():
    try:
        file = save_file(request.files["file"])
        target_lang = request.form.get("lang", "es") # Default Spanish
        output = file.replace(".pdf", f"_translated_{target_lang}.txt")
        
        doc = fitz.open(file)
        full_text = ""
        for page in doc:
            full_text += clean_extracted_text(page.get_text())
            
        # Translate chunks (limit 5000 chars roughly)
        translator = GoogleTranslator(source='auto', target=target_lang)
        translated_text = ""
        
        # Simple chunking
        chunk_size = 4500
        for i in range(0, len(full_text), chunk_size):
            chunk = full_text[i:i+chunk_size]
            if chunk.strip():
                try:
                    translated_text += translator.translate(chunk) + "\n\n"
                except Exception as tr_err:
                    print(f"Translation chunk error: {tr_err}")
                    translated_text += f"[Translation Error Block]\n{chunk}\n\n"
        
        # Output as a real PDF using reportlab instead of .txt
        output = file.replace(".pdf", f"_translated_{target_lang}.pdf")
        doc_pdf = SimpleDocTemplate(output, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Split text by newlines and add as paragraphs to handle wrapping
        for para_text in translated_text.split("\n\n"):
            if para_text.strip():
                clean_para = para_text.replace("\n", "<br/>")
                story.append(Paragraph(clean_para, styles["Normal"]))
                story.append(Spacer(1, 12))
                
        doc_pdf.build(story)

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# AI SUMMARIZER
# ------------------------
@app.route("/summarize", methods=["POST"])
def summarize():
    try:
        file = save_file(request.files["file"])
        output = file.replace(".pdf", "_summary.txt")
        
        # Extract text
        doc = fitz.open(file)
        text = ""
        for page in doc:
            page_text = clean_extracted_text(page.get_text())
            text += page_text
            
        # Basic Extractive Summary (Frequency based)
        import heapq
        import re
        
        # Remove citations, brackets
        clean_text = re.sub(r'\[[0-9]*\]', ' ', text)
        clean_text = re.sub(r'\s+', ' ', clean_text)
        
        formatted_text = re.sub('[^a-zA-Z]', ' ', clean_text)
        formatted_text = re.sub(r'\s+', ' ', formatted_text)
        
        sentence_list = re.split(r'(?<=[.!?])\s+', clean_text)
        
        stopwords = ["the", "and", "of", "to", "in", "a", "is", "that", "for", "it", "on", "as", "are", "with", "this", "was", "or", "an", "be", "by", "not", "at", "from"]
        
        word_frequencies = {}
        for word in formatted_text.lower().split():
            if word not in stopwords:
                if word not in word_frequencies.keys():
                    word_frequencies[word] = 1
                else:
                    word_frequencies[word] += 1
                    
        if not word_frequencies:
             return jsonify({"error": "Not enough text to summarize"}), 400

        maximum_frequncy = max(word_frequencies.values())
        for word in word_frequencies.keys():
            word_frequencies[word] = (word_frequencies[word]/maximum_frequncy)
            
        sentence_scores = {}
        for sent in sentence_list:
            for word in sent.lower().split():
                if word in word_frequencies.keys():
                    if len(sent.split(' ')) < 30: # limit long sentences
                        if sent not in sentence_scores.keys():
                            sentence_scores[sent] = word_frequencies[word]
                        else:
                            sentence_scores[sent] += word_frequencies[word]
                            
        summary_sentences = heapq.nlargest(7, sentence_scores, key=sentence_scores.get)
        summary = ' '.join(summary_sentences)
        
        with open(output, "w", encoding="utf-8-sig") as f:
            f.write("--- SUMMARY ---\n\n")
            f.write(summary)
            f.write("\n\n--- ORIGINAL TEXT START ---\n")
            f.write(text[:1000] + "\n...") # Preview

        auto_delete(file)
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# WORKFLOW (Chain)
# ------------------------
@app.route("/workflow", methods=["POST"])
def workflow():
    try:
        # Receives file + JSON actions
        # actions = [{"tool": "rotate", "params": {"angle": 90}}, {"tool": "compress", "params": {}}]
        file = save_file(request.files["file"])
        actions_str = request.form.get("actions", "[]")
        try:
             import json
             actions = json.loads(actions_str)
        except:
             return jsonify({"error": "Invalid actions JSON"}), 400
             
        current_file = file
        temp_files = [file]
        
        # Map tool names to internal functions or logic
        # Since logic is inside route functions which return responses, we can't easily call them directly without mocking request/context.
        # So we have to allow some logic reuse or internal implementations.
        # Refactoring everything to separate logic functions would be best.
        # For now, I'll allow simple operations that I can re-implement or call logic for.
        
        reader = PdfReader(current_file)
        writer = PdfWriter()
        writer.append(reader)
        
        # Simple memory-based processing for basics
        for action in actions:
            tool = action.get("tool")
            params = action.get("params", {})
            
            # Re-read writer to reader if we modified it?
            # Better: Write to temp buffer, read back
            out_buffer = io.BytesIO()
            writer.write(out_buffer)
            out_buffer.seek(0)
            
            # Apply Action
            if tool == "rotate":
                angle = int(params.get("angle", 90))
                reader_step = PdfReader(out_buffer)
                writer = PdfWriter()
                for page in reader_step.pages:
                    page.rotate(angle)
                    writer.add_page(page)
                    
            elif tool == "remove-pages":
                # params: pages "1,2"
                pages_to_remove = params.get("pages", "")
                reader_step = PdfReader(out_buffer)
                writer = PdfWriter()
                total = len(reader_step.pages)
                
                to_remove = set()
                for part in pages_to_remove.split(','):
                    if '-' in part:
                        s, e = map(int, part.split('-'))
                        to_remove.update(range(s, e + 1))
                    else:
                        to_remove.add(int(part))
                        
                for i in range(total):
                    if (i+1) not in to_remove:
                        writer.add_page(reader_step.pages[i])

            elif tool == "compress":
                 # Metadata compression only in pypdf
                 reader_step = PdfReader(out_buffer)
                 writer = PdfWriter()
                 writer.append(reader_step)
                 writer.compress_identical_objects()
            
            # Add more supported workflow actions here
            
        output = os.path.join(UPLOAD_FOLDER, f"workflow_{uuid.uuid4()}.pdf")
        with open(output, "wb") as f:
            writer.write(f)
            
        # Clean inputs
        for p in temp_files:
            if os.path.exists(p):
                try: 
                     # auto_delete handles it, or we delete here
                     pass
                except: pass
                
        auto_delete(output)
        return send_file(output, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# GROUP 5: EXTRACTION & UTILITIES
# ------------------------

# ------------------------
# PDF EXTRACT
# ------------------------
@app.route("/pdf-extract", methods=["POST"])
def pdf_extract():
    try:
        file = save_file(request.files["file"])
        base_name = os.path.splitext(os.path.basename(file))[0]
        extract_dir = os.path.join(UPLOAD_FOLDER, f"extract_{uuid.uuid4()}")
        os.makedirs(extract_dir, exist_ok=True)
        
        doc = fitz.open(file)
        
        # 1. Metadata
        import json
        metadata_path = os.path.join(extract_dir, "metadata.json")
        with open(metadata_path, "w", encoding="utf-8-sig") as f:
            json.dump(doc.metadata, f, indent=4)
            
        # 2. Text
        text_path = os.path.join(extract_dir, "content.txt")
        full_text = ""
        for page in doc:
            full_text += clean_extracted_text(page.get_text())
        with open(text_path, "w", encoding="utf-8-sig") as f:
            f.write(full_text)
            
        # 3. Images
        images_dir = os.path.join(extract_dir, "images")
        os.makedirs(images_dir, exist_ok=True)
        
        for i, page in enumerate(doc):
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = f"page{i+1}_img{img_index+1}.{image_ext}"
                with open(os.path.join(images_dir, image_filename), "wb") as f:
                    f.write(image_bytes)
                    
        # Zip it up
        output_zip = os.path.join(UPLOAD_FOLDER, f"{base_name}_extracted.zip")
        import shutil
        shutil.make_archive(output_zip.replace(".zip", ""), 'zip', extract_dir)
        
        # Cleanup extract dir
        shutil.rmtree(extract_dir)
        auto_delete(file)
        auto_delete(output_zip)
        
        return send_file(output_zip, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# COPY PDF
# ------------------------
@app.route("/copy-pdf", methods=["POST"])
def copy_pdf():
    try:
        file = save_file(request.files["file"])
        copies = int(request.form.get("copies", 1))
        
        if copies < 1:
            copies = 1
            
        # If copies = 1, just return the file? 
        # Typically "Copy" tool means duplicate, so maybe return a zip if > 1, or just download the same file renamed?
        # Let's support generating a ZIP of N copies.
        
        output_zip = os.path.join(UPLOAD_FOLDER, f"copies_{uuid.uuid4()}.zip")
        import zipfile
        
        with zipfile.ZipFile(output_zip, 'w') as zipf:
            base_name = os.path.basename(file)
            name, ext = os.path.splitext(base_name)
            for i in range(copies):
                # Copy file logic
                # Write the same source file to zip with different names
                zipf.write(file, arcname=f"{name}_copy_{i+1}{ext}")
                
        auto_delete(file)
        auto_delete(output_zip)
        return send_file(output_zip, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ------------------------
# RUN
# ------------------------

if __name__ == "__main__":
    app.run(host="0.0.0.0", debug=False)
